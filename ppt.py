from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color Palette - Forest & Tech theme
C_DARK = RGBColor(0x0D, 0x2B, 0x1E)       # Very dark green (bg for title slides)
C_GREEN = RGBColor(0x1A, 0x73, 0x48)      # Forest green (primary)
C_MINT = RGBColor(0x2E, 0xCC, 0x71)       # Mint green (accent)
C_LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)# Light green (card bg)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT = RGBColor(0x1A, 0x25, 0x1E)  # Very dark for body text
C_GRAY = RGBColor(0x5D, 0x6D, 0x7E)
C_LIGHT_BG = RGBColor(0xF0, 0xFB, 0xF4)  # Very light green bg
C_ACCENT = RGBColor(0xFF, 0x8C, 0x00)     # Orange accent
C_RED = RGBColor(0xE7, 0x4C, 0x3C)
C_BLUE = RGBColor(0x21, 0x80, 0xC8)

def hex_to_rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, x, y, w, h, color, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if transparency > 0:
        shape.fill.fore_color.theme_color = None
    return shape

def add_text(slide, text, x, y, w, h, font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox

def add_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_card(slide, x, y, w, h, bg_color=None):
    if bg_color is None:
        bg_color = C_WHITE
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xDD, 0xEE, 0xE5)
    shape.line.width = Pt(0.75)
    return shape

def add_accent_bar(slide, x, y, h, color=None):
    if color is None:
        color = C_MINT
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

# ─────────────────────────────────────────────
# SLIDE 1: Title Slide
# ─────────────────────────────────────────────
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide1, C_DARK)

# Decorative green band on left
left_band = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
left_band.fill.solid(); left_band.fill.fore_color.rgb = C_GREEN; left_band.line.fill.background()

# Accent bottom band
bot = slide1.shapes.add_shape(1, Inches(0), Inches(6.8), Inches(13.33), Inches(0.7))
bot.fill.solid(); bot.fill.fore_color.rgb = C_GREEN; bot.line.fill.background()

# Decorative circle top-right
cr = slide1.shapes.add_shape(9, Inches(10.5), Inches(-0.5), Inches(3), Inches(3))
cr.fill.solid(); cr.fill.fore_color.rgb = RGBColor(0x1A,0x73,0x48)
cr.fill.fore_color.rgb = C_GREEN; cr.line.fill.background()

cr2 = slide1.shapes.add_shape(9, Inches(11), Inches(0), Inches(2), Inches(2))
cr2.fill.solid(); cr2.fill.fore_color.rgb = C_MINT; cr2.line.fill.background()

# Badge
badge = slide1.shapes.add_shape(1, Inches(0.9), Inches(1.2), Inches(2.8), Inches(0.45))
badge.fill.solid(); badge.fill.fore_color.rgb = C_MINT; badge.line.fill.background()
txb = slide1.shapes.add_textbox(Inches(0.9), Inches(1.2), Inches(2.8), Inches(0.45))
txb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
r = txb.text_frame.paragraphs[0].add_run()
r.text = "MACHINE LEARNING PROJECT"
r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = C_DARK

# Title
t = slide1.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(9), Inches(1.5))
t.text_frame.word_wrap = True
p = t.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "Intelligent Irrigation"
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = C_WHITE

t2 = slide1.shapes.add_textbox(Inches(0.9), Inches(3.3), Inches(9), Inches(0.9))
p2 = t2.text_frame.paragraphs[0]
r2 = p2.add_run(); r2.text = "System Using Machine Learning"
r2.font.size = Pt(36); r2.font.bold = True; r2.font.color.rgb = C_MINT

# Subtitle
add_text(slide1, "A Data-Driven Approach to Smart Water Management & Agricultural Optimization",
         0.9, 4.3, 9, 0.7, font_size=14, color=RGBColor(0xB2,0xD8,0xC4), italic=True)

# Stats row at bottom
for i, (val, lbl) in enumerate([("4", "Input\nParameters"), ("6", "System\nModules"), ("ML", "Core\nEngine"), ("Real-Time", "Monitoring")]):
    bx = 0.9 + i * 2.9
    sq = slide1.shapes.add_shape(1, Inches(bx), Inches(5.4), Inches(2.5), Inches(1.1))
    sq.fill.solid(); sq.fill.fore_color.rgb = RGBColor(0x1E,0x3D,0x2D); sq.line.fill.background()
    add_text(slide1, val, bx+0.08, 5.42, 2.3, 0.5, font_size=22, bold=True, color=C_MINT, align=PP_ALIGN.CENTER)
    add_text(slide1, lbl, bx+0.08, 5.85, 2.3, 0.45, font_size=9, color=RGBColor(0xB2,0xD8,0xC4), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 2: Problem Statement
# ─────────────────────────────────────────────
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide2, C_LIGHT_BG)

# Header band
hdr = slide2.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
hdr.fill.solid(); hdr.fill.fore_color.rgb = C_GREEN; hdr.line.fill.background()
add_text(slide2, "Problem Statement", 0.4, 0.15, 10, 0.8, font_size=28, bold=True, color=C_WHITE)
add_text(slide2, "Slide 02", 11.8, 0.25, 1.3, 0.5, font_size=11, color=RGBColor(0xB2,0xD8,0xC4), align=PP_ALIGN.RIGHT)

problems = [
    ("💧", "Manual Irrigation Methods", "Traditional methods rely on farmer intuition, fixed schedules, and manual observation — leading to inconsistent water application and poor crop outcomes."),
    ("🌊", "Water Wastage", "Over-irrigation is widespread. Excess water application depletes groundwater, increases runoff, and results in significant economic losses for farmers."),
    ("📡", "Lack of Real-Time Monitoring", "No continuous sensor-based feedback means irrigation decisions are made without awareness of current soil moisture, weather, or crop conditions."),
    ("⚡", "High Electricity Usage", "Inefficient pump operation schedules — running motors during unnecessary periods — significantly inflate electricity bills and carbon emissions."),
    ("📊", "Absence of Data-Driven Decisions", "Without historical data analysis or predictive models, farmers cannot anticipate irrigation needs, preventing proactive and optimized farm management."),
]

cols = [(0.35, 5.9), (6.85, 5.9)]
positions = [(0.35, 1.3), (6.85, 1.3), (0.35, 3.35), (6.85, 3.35), (3.6, 5.4)]

for idx, (icon, title, desc) in enumerate(problems):
    if idx < 4:
        cx, _ = positions[idx]
        cy = positions[idx][1]
        cw = 5.9
    else:
        cx, cy = 3.6, 5.4
        cw = 5.9

    card = slide2.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(cw), Inches(1.85) if idx < 4 else Inches(1.7))
    card.fill.solid(); card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = RGBColor(0xC0,0xDF,0xCC); card.line.width = Pt(1)

    # Accent bar
    acc = slide2.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(0.06), Inches(1.85) if idx < 4 else Inches(1.7))
    acc.fill.solid(); acc.fill.fore_color.rgb = C_MINT; acc.line.fill.background()

    add_text(slide2, f"{icon} {title}", cx+0.15, cy+0.1, cw-0.2, 0.45, font_size=13, bold=True, color=C_GREEN)
    add_text(slide2, desc, cx+0.15, cy+0.52, cw-0.25, 1.2, font_size=10.5, color=C_DARK_TEXT)

# ─────────────────────────────────────────────
# SLIDE 3: Objective
# ─────────────────────────────────────────────
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide3, C_LIGHT_BG)

hdr = slide3.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
hdr.fill.solid(); hdr.fill.fore_color.rgb = C_GREEN; hdr.line.fill.background()
add_text(slide3, "Project Objective", 0.4, 0.15, 10, 0.8, font_size=28, bold=True, color=C_WHITE)
add_text(slide3, "Slide 03", 11.8, 0.25, 1.3, 0.5, font_size=11, color=RGBColor(0xB2,0xD8,0xC4), align=PP_ALIGN.RIGHT)

# Big objective statement
stmt = slide3.shapes.add_shape(1, Inches(0.4), Inches(1.2), Inches(12.5), Inches(1.15))
stmt.fill.solid(); stmt.fill.fore_color.rgb = C_GREEN; stmt.line.fill.background()
add_text(slide3, "\"Develop an Intelligent, ML-Powered Irrigation System that Predicts Water Requirements and Optimizes Agricultural Efficiency\"",
         0.55, 1.3, 12.2, 0.9, font_size=13.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, italic=True)

objectives = [
    ("01", "Intelligent Irrigation System", "Build a smart, automated irrigation controller that adapts to real-world environmental conditions rather than fixed schedules."),
    ("02", "Machine Learning Techniques", "Apply supervised learning algorithms to analyze patterns in historical data and make data-driven irrigation decisions."),
    ("03", "Predict Irrigation Requirements", "Forecast when and how much water crops need based on sensor inputs, eliminating over- and under-watering."),
    ("04", "Optimize Water Usage & Efficiency", "Reduce water consumption by 30–50% through precision irrigation, lowering costs and environmental impact."),
]

for idx, (num, title, desc) in enumerate(objectives):
    col = idx % 2
    row = idx // 2
    cx = 0.4 + col * 6.5
    cy = 2.55 + row * 2.3

    card = slide3.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(6.2), Inches(2.0))
    card.fill.solid(); card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = RGBColor(0xC0,0xDF,0xCC); card.line.width = Pt(1)

    num_circle = slide3.shapes.add_shape(9, Inches(cx+0.15), Inches(cy+0.2), Inches(0.6), Inches(0.6))
    num_circle.fill.solid(); num_circle.fill.fore_color.rgb = C_MINT; num_circle.line.fill.background()
    add_text(slide3, num, cx+0.15, cy+0.22, 0.6, 0.5, font_size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    add_text(slide3, title, cx+0.9, cy+0.18, 5.1, 0.5, font_size=13, bold=True, color=C_GREEN)
    add_text(slide3, desc, cx+0.15, cy+0.88, 5.85, 1.0, font_size=10.5, color=C_DARK_TEXT)

# ─────────────────────────────────────────────
# SLIDE 4: Proposed Solution
# ─────────────────────────────────────────────
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide4, C_LIGHT_BG)

hdr = slide4.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
hdr.fill.solid(); hdr.fill.fore_color.rgb = C_GREEN; hdr.line.fill.background()
add_text(slide4, "Proposed Solution", 0.4, 0.15, 10, 0.8, font_size=28, bold=True, color=C_WHITE)
add_text(slide4, "Slide 04", 11.8, 0.25, 1.3, 0.5, font_size=11, color=RGBColor(0xB2,0xD8,0xC4), align=PP_ALIGN.RIGHT)

# Left column - solution description
left_bg = slide4.shapes.add_shape(1, Inches(0.35), Inches(1.25), Inches(5.8), Inches(5.9))
left_bg.fill.solid(); left_bg.fill.fore_color.rgb = C_WHITE
left_bg.line.color.rgb = RGBColor(0xC0,0xDF,0xCC); left_bg.line.width = Pt(1)

acc = slide4.shapes.add_shape(1, Inches(0.35), Inches(1.25), Inches(0.07), Inches(5.9))
acc.fill.solid(); acc.fill.fore_color.rgb = C_GREEN; acc.line.fill.background()

add_text(slide4, "Our Solution", 0.55, 1.35, 5.4, 0.55, font_size=18, bold=True, color=C_GREEN)
add_text(slide4, "A Data-Driven Intelligent Irrigation System", 0.55, 1.85, 5.4, 0.5, font_size=13, bold=True, color=C_DARK_TEXT)

sol_points = [
    "Collects real-time environmental data from multiple sensors",
    "Processes and analyzes data using ML algorithms",
    "Predicts irrigation ON/OFF status with high accuracy",
    "Provides smart, actionable recommendations to farmers",
    "Continuously learns and improves from new data",
    "Reduces water waste and optimizes energy usage",
]
for i, pt in enumerate(sol_points):
    dot = slide4.shapes.add_shape(9, Inches(0.55), Inches(2.55+i*0.52), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = C_MINT; dot.line.fill.background()
    add_text(slide4, pt, 0.85, 2.5+i*0.52, 5.1, 0.45, font_size=11, color=C_DARK_TEXT)

# Right column - key features
features = [
    ("🔢", "Data-Driven", "Uses real sensor data, not guesswork"),
    ("🤖", "ML-Powered", "Supervised learning at its core"),
    ("⚡", "Real-Time", "Instant irrigation decisions"),
    ("💡", "Smart Alerts", "Actionable recommendations"),
]

for i, (icon, title, desc) in enumerate(features):
    row = i // 2; col = i % 2
    cx = 6.5 + col * 3.35
    cy = 1.3 + row * 2.7
    card = slide4.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(3.1), Inches(2.4))
    card.fill.solid()
    card.fill.fore_color.rgb = C_GREEN if i % 2 == 0 else RGBColor(0x14,0x5A,0x32)
    card.line.fill.background()

    add_text(slide4, icon, cx+0.15, cy+0.2, 0.7, 0.6, font_size=24, align=PP_ALIGN.CENTER)
    add_text(slide4, title, cx+0.15, cy+0.85, 2.7, 0.5, font_size=16, bold=True, color=C_MINT)
    add_text(slide4, desc, cx+0.15, cy+1.35, 2.8, 0.8, font_size=11, color=RGBColor(0xCC,0xEE,0xDD))

# ─────────────────────────────────────────────
# SLIDE 5: Input Parameters
# ─────────────────────────────────────────────
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide5, C_LIGHT_BG)

hdr = slide5.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
hdr.fill.solid(); hdr.fill.fore_color.rgb = C_GREEN; hdr.line.fill.background()
add_text(slide5, "Input Parameters", 0.4, 0.15, 10, 0.8, font_size=28, bold=True, color=C_WHITE)
add_text(slide5, "Slide 05", 11.8, 0.25, 1.3, 0.5, font_size=11, color=RGBColor(0xB2,0xD8,0xC4), align=PP_ALIGN.RIGHT)

params = [
    ("💧", "Soil Moisture", "Measures water content in soil (%)", "Critical indicator", "Determines when soil is too dry or saturated for optimal crop growth", RGBColor(0x21,0x80,0xC8)),
    ("🌡️", "Temperature", "Ambient air temperature (°C)", "Evaporation driver", "High temperatures increase evapotranspiration, requiring more frequent irrigation", RGBColor(0xE7,0x4C,0x3C)),
    ("💨", "Humidity", "Relative humidity of air (%)", "Moisture balance", "Low humidity increases water loss from soil surface and plant transpiration rate", RGBColor(0x9B,0x59,0xB6)),
    ("🌧️", "Rainfall", "Precipitation amount (mm)", "Natural supply", "Recent rainfall reduces irrigation needs significantly; prevents unnecessary watering", RGBColor(0x1A,0xBC,0x9C)),
]

for idx, (icon, name, measure, badge_text, detail, color) in enumerate(params):
    cx = 0.35 + idx * 3.2
    card = slide5.shapes.add_shape(1, Inches(cx), Inches(1.3), Inches(2.95), Inches(5.85))
    card.fill.solid(); card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = color; card.line.width = Pt(1.5)

    top = slide5.shapes.add_shape(1, Inches(cx), Inches(1.3), Inches(2.95), Inches(1.6))
    top.fill.solid(); top.fill.fore_color.rgb = color; top.line.fill.background()

    add_text(slide5, icon, cx+0.1, 1.35, 2.75, 0.75, font_size=32, align=PP_ALIGN.CENTER)
    add_text(slide5, name, cx+0.1, 2.05, 2.75, 0.65, font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide5, measure, cx+0.12, 3.05, 2.7, 0.5, font_size=11, bold=True, color=color)

    badge = slide5.shapes.add_shape(1, Inches(cx+0.12), Inches(3.65), Inches(2.7), Inches(0.35))
    badge.fill.solid(); badge.fill.fore_color.rgb = RGBColor(0xE8,0xF8,0xF5); badge.line.fill.background()
    add_text(slide5, f"▶  {badge_text}", cx+0.18, 3.66, 2.6, 0.33, font_size=10, bold=True, color=color)

    add_text(slide5, detail, cx+0.12, 4.15, 2.7, 1.75, font_size=10.5, color=C_DARK_TEXT)

# ─────────────────────────────────────────────
# SLIDE 6: System Workflow
# ─────────────────────────────────────────────
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide6, C_LIGHT_BG)

hdr = slide6.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
hdr.fill.solid(); hdr.fill.fore_color.rgb = C_GREEN; hdr.line.fill.background()
add_text(slide6, "System Workflow", 0.4, 0.15, 10, 0.8, font_size=28, bold=True, color=C_WHITE)
add_text(slide6, "Slide 06", 11.8, 0.25, 1.3, 0.5, font_size=11, color=RGBColor(0xB2,0xD8,0xC4), align=PP_ALIGN.RIGHT)

steps = [
    ("01", "Data\nCollection", "Sensors gather soil moisture, temperature, humidity & rainfall data in real-time", "📥", C_BLUE),
    ("02", "Data\nPreprocessing", "Clean, normalize, and handle missing values; feature engineering for ML readiness", "⚙️", RGBColor(0x8E,0x44,0xAD)),
    ("03", "EDA", "Visualize distributions, correlations, and patterns in environmental data", "🔍", RGBColor(0xE6,0x7E,0x22)),
    ("04", "Model\nTraining", "Train supervised ML classifiers on labeled irrigation dataset with cross-validation", "🤖", C_GREEN),
    ("05", "Prediction", "Real-time inference: predict Irrigation ON/OFF status for current conditions", "⚡", RGBColor(0xE7,0x4C,0x3C)),
    ("06", "Visualization", "Display results, trends, accuracy metrics, and recommendations via dashboard", "📊", RGBColor(0x1A,0xBC,0x9C)),
]

for idx, (num, title, desc, icon, color) in enumerate(steps):
    col = idx % 3
    row = idx // 2
    cx = 0.35 + col * 4.35
    cy = 1.3 + row * 2.85

    if idx < 3:
        cy = 1.3
    else:
        cy = 4.15

    card = slide6.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(4.1), Inches(2.65))
    card.fill.solid(); card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = color; card.line.width = Pt(1.5)

    top_bar = slide6.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(4.1), Inches(0.06))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.fill.background()

    num_c = slide6.shapes.add_shape(9, Inches(cx+0.15), Inches(cy+0.18), Inches(0.55), Inches(0.55))
    num_c.fill.solid(); num_c.fill.fore_color.rgb = color; num_c.line.fill.background()
    add_text(slide6, num, cx+0.15, cy+0.2, 0.55, 0.45, font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide6, icon, cx+3.35, cy+0.12, 0.6, 0.55, font_size=22, align=PP_ALIGN.CENTER)
    add_text(slide6, title, cx+0.85, cy+0.18, 2.4, 0.65, font_size=13, bold=True, color=color)
    add_text(slide6, desc, cx+0.15, cy+0.95, 3.8, 1.5, font_size=10.5, color=C_DARK_TEXT)

    # Arrow between steps in same row
    if idx < 2 or (idx > 2 and idx < 5):
        arr = slide6.shapes.add_shape(1, Inches(cx+4.1), Inches(cy+1.05), Inches(0.25), Inches(0.4))
        arr.fill.solid(); arr.fill.fore_color.rgb = color; arr.line.fill.background()
        add_text(slide6, "▶", cx+4.1, cy+1.08, 0.25, 0.35, font_size=14, color=color, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 7: System Architecture
# ─────────────────────────────────────────────
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide7, C_LIGHT_BG)

hdr = slide7.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
hdr.fill.solid(); hdr.fill.fore_color.rgb = C_GREEN; hdr.line.fill.background()
add_text(slide7, "System Architecture", 0.4, 0.15, 10, 0.8, font_size=28, bold=True, color=C_WHITE)
add_text(slide7, "Slide 07", 11.8, 0.25, 1.3, 0.5, font_size=11, color=RGBColor(0xB2,0xD8,0xC4), align=PP_ALIGN.RIGHT)

# Three main layers: Input -> Processing -> Output
layers = [
    ("INPUT LAYER", C_BLUE, ["Soil Moisture\nSensor", "Temperature\nSensor", "Humidity\nSensor", "Rainfall\nSensor"], 0.35),
    ("PROCESSING LAYER", C_GREEN, ["Data\nCleaning", "Feature\nEngineering", "ML Model\n(Core)", "Model\nEvaluation"], 4.75),
    ("OUTPUT LAYER", RGBColor(0xE6,0x7E,0x22), ["Irrigation\nON/OFF", "Smart\nAlert", "Dashboard\nReport", "Next-Day\nForecast"], 9.15),
]

for lx, (layer_name, lcolor, items, sx) in enumerate(layers):
    # Layer header
    lhdr = slide7.shapes.add_shape(1, Inches(sx), Inches(1.2), Inches(3.85), Inches(0.5))
    lhdr.fill.solid(); lhdr.fill.fore_color.rgb = lcolor; lhdr.line.fill.background()
    add_text(slide7, layer_name, sx+0.05, 1.23, 3.75, 0.42, font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Layer bg
    lbg = slide7.shapes.add_shape(1, Inches(sx), Inches(1.7), Inches(3.85), Inches(4.95))
    lbg.fill.solid(); lbg.fill.fore_color.rgb = C_WHITE
    lbg.line.color.rgb = lcolor; lbg.line.width = Pt(1.5)

    for i, item_text in enumerate(items):
        iy = 1.9 + i * 1.15
        ic = slide7.shapes.add_shape(1, Inches(sx+0.3), Inches(iy), Inches(3.25), Inches(0.85))
        ic.fill.solid()
        ic.fill.fore_color.rgb = RGBColor(0xE8,0xF5,0xFE) if lx == 0 else (C_LIGHT_GREEN if lx == 1 else RGBColor(0xFE,0xF5,0xE8))
        ic.line.color.rgb = lcolor; ic.line.width = Pt(0.75)
        add_text(slide7, item_text, sx+0.35, iy+0.1, 3.15, 0.65, font_size=11, bold=(i==2 and lx==1), color=lcolor if (i==2 and lx==1) else C_DARK_TEXT, align=PP_ALIGN.CENTER)

    # Arrow to next layer
    if lx < 2:
        arr_x = sx + 3.85
        arr_bg = slide7.shapes.add_shape(1, Inches(arr_x), Inches(3.6), Inches(0.55), Inches(0.6))
        arr_bg.fill.solid(); arr_bg.fill.fore_color.rgb = lcolor; arr_bg.line.fill.background()
        add_text(slide7, "▶", arr_x+0.05, 3.63, 0.45, 0.5, font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ML Core highlight
ml_note = slide7.shapes.add_shape(1, Inches(4.75), Inches(6.75), Inches(3.85), Inches(0.45))
ml_note.fill.solid(); ml_note.fill.fore_color.rgb = C_GREEN; ml_note.line.fill.background()
add_text(slide7, "⭐ ML Model is the decision-making core of the entire architecture", 4.8, 6.77, 3.75, 0.4, font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 8: Machine Learning Model
# ─────────────────────────────────────────────
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide8, C_LIGHT_BG)

hdr = slide8.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
hdr.fill.solid(); hdr.fill.fore_color.rgb = C_GREEN; hdr.line.fill.background()
add_text(slide8, "Machine Learning Model", 0.4, 0.15, 10, 0.8, font_size=28, bold=True, color=C_WHITE)
add_text(slide8, "Slide 08", 11.8, 0.25, 1.3, 0.5, font_size=11, color=RGBColor(0xB2,0xD8,0xC4), align=PP_ALIGN.RIGHT)

# Left - ML info
left_card = slide8.shapes.add_shape(1, Inches(0.35), Inches(1.3), Inches(6.0), Inches(5.9))
left_card.fill.solid(); left_card.fill.fore_color.rgb = C_WHITE
left_card.line.color.rgb = RGBColor(0xC0,0xDF,0xCC); left_card.line.width = Pt(1)

add_text(slide8, "Model Overview", 0.55, 1.4, 5.6, 0.5, font_size=16, bold=True, color=C_GREEN)

ml_info = [
    ("Approach", "Supervised Classification Learning"),
    ("Algorithm", "Decision Tree / Random Forest / SVM"),
    ("Target Variable", "Irrigation Status (ON = 1 / OFF = 0)"),
    ("Training Data", "Historical labeled environmental records"),
    ("Validation", "Train-Test Split (80:20 ratio)"),
    ("Metrics", "Accuracy, Precision, Recall, F1-Score"),
]

for i, (lbl, val) in enumerate(ml_info):
    ry = 2.05 + i * 0.8
    row_bg = slide8.shapes.add_shape(1, Inches(0.45), Inches(ry), Inches(5.7), Inches(0.62))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = C_LIGHT_GREEN if i % 2 == 0 else C_WHITE
    row_bg.line.fill.background()

    add_text(slide8, lbl, 0.55, ry+0.08, 2.2, 0.45, font_size=11, bold=True, color=C_GREEN)
    add_text(slide8, val, 2.8, ry+0.08, 3.2, 0.45, font_size=11, color=C_DARK_TEXT)

# Right - Metrics visual
right_card = slide8.shapes.add_shape(1, Inches(6.7), Inches(1.3), Inches(6.2), Inches(5.9))
right_card.fill.solid(); right_card.fill.fore_color.rgb = C_WHITE
right_card.line.color.rgb = RGBColor(0xC0,0xDF,0xCC); right_card.line.width = Pt(1)

add_text(slide8, "Evaluation Metrics & Workflow", 6.9, 1.4, 5.8, 0.5, font_size=16, bold=True, color=C_GREEN)

metrics = [("Accuracy", "95%+", C_MINT), ("Precision", "93%+", C_BLUE), ("Recall", "94%+", RGBColor(0xE6,0x7E,0x22)), ("F1-Score", "93%+", RGBColor(0x9B,0x59,0xB6))]
for i, (metric, val, mc) in enumerate(metrics):
    mx = 6.9 + (i % 2) * 2.9
    my = 2.05 + (i // 2) * 1.5
    mc_card = slide8.shapes.add_shape(1, Inches(mx), Inches(my), Inches(2.6), Inches(1.2))
    mc_card.fill.solid(); mc_card.fill.fore_color.rgb = mc; mc_card.line.fill.background()
    add_text(slide8, val, mx+0.1, my+0.1, 2.4, 0.65, font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide8, metric, mx+0.1, my+0.75, 2.4, 0.35, font_size=12, color=RGBColor(0xDD,0xFF,0xEE), align=PP_ALIGN.CENTER)

# Training flow
tf_items = ["Collect Data", "→", "Preprocess", "→", "Train Model", "→", "Evaluate", "→", "Deploy"]
add_text(slide8, "Model Training Flow:", 6.9, 5.2, 5.8, 0.4, font_size=12, bold=True, color=C_GREEN)

flow_bg = slide8.shapes.add_shape(1, Inches(6.75), Inches(5.65), Inches(5.9), Inches(1.3))
flow_bg.fill.solid(); flow_bg.fill.fore_color.rgb = C_LIGHT_GREEN; flow_bg.line.fill.background()

flow_text = " → ".join(["Collect Data", "Preprocess", "Train", "Evaluate", "Deploy"])
add_text(slide8, flow_text, 6.85, 5.95, 5.7, 0.7, font_size=11, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 9: System Modules
# ─────────────────────────────────────────────
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide9, C_LIGHT_BG)

hdr = slide9.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
hdr.fill.solid(); hdr.fill.fore_color.rgb = C_GREEN; hdr.line.fill.background()
add_text(slide9, "System Modules", 0.4, 0.15, 10, 0.8, font_size=28, bold=True, color=C_WHITE)
add_text(slide9, "Slide 09", 11.8, 0.25, 1.3, 0.5, font_size=11, color=RGBColor(0xB2,0xD8,0xC4), align=PP_ALIGN.RIGHT)

modules = [
    ("📥", "Data Collection", C_BLUE,
     ["Real-time sensor data ingestion", "Supports multiple sensor types", "Timestamped data logging", "Data buffering & storage"]),
    ("⚙️", "Data Preprocessing", RGBColor(0x8E,0x44,0xAD),
     ["Missing value imputation", "Outlier detection & removal", "Feature normalization (0-1)", "Train/test split generation"]),
    ("🔍", "Data Analysis (EDA)", RGBColor(0xE6,0x7E,0x22),
     ["Statistical summaries", "Correlation heatmaps", "Distribution plots", "Seasonal trend analysis"]),
    ("🤖", "ML Model Module", C_GREEN,
     ["Model selection & training", "Hyperparameter tuning", "Cross-validation", "Model persistence (save/load)"]),
    ("⚡", "Prediction Engine", RGBColor(0xE7,0x4C,0x3C),
     ["Real-time inference", "ON/OFF classification", "Confidence scoring", "Threshold optimization"]),
    ("📊", "Visualization", RGBColor(0x1A,0xBC,0x9C),
     ["Interactive dashboards", "Historical trend charts", "Accuracy metric display", "Alert & notification UI"]),
]

for idx, (icon, title, color, points) in enumerate(modules):
    col = idx % 3
    row = idx // 3
    cx = 0.35 + col * 4.35
    cy = 1.3 + row * 3.1

    card = slide9.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(4.1), Inches(2.85))
    card.fill.solid(); card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = color; card.line.width = Pt(1.5)

    top_h = slide9.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(4.1), Inches(0.75))
    top_h.fill.solid(); top_h.fill.fore_color.rgb = color; top_h.line.fill.background()

    add_text(slide9, icon, cx+0.1, cy+0.05, 0.6, 0.6, font_size=18, align=PP_ALIGN.CENTER)
    add_text(slide9, title, cx+0.75, cy+0.1, 3.2, 0.55, font_size=14, bold=True, color=C_WHITE)

    for pi, pt in enumerate(points):
        py = cy + 0.88 + pi * 0.47
        dot = slide9.shapes.add_shape(9, Inches(cx+0.18), Inches(py+0.1), Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
        add_text(slide9, pt, cx+0.4, py+0.02, 3.55, 0.4, font_size=10, color=C_DARK_TEXT)

# ─────────────────────────────────────────────
# SLIDE 10: Activity Diagram / Flowchart
# ─────────────────────────────────────────────
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide10, C_LIGHT_BG)

hdr = slide10.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
hdr.fill.solid(); hdr.fill.fore_color.rgb = C_GREEN; hdr.line.fill.background()
add_text(slide10, "Activity Diagram — Irrigation System Workflow", 0.4, 0.15, 10.5, 0.8, font_size=24, bold=True, color=C_WHITE)
add_text(slide10, "Slide 10", 11.8, 0.25, 1.3, 0.5, font_size=11, color=RGBColor(0xB2,0xD8,0xC4), align=PP_ALIGN.RIGHT)

# Flowchart - vertical flow with decision diamond
# START
start_circle = slide10.shapes.add_shape(9, Inches(5.9), Inches(1.2), Inches(1.0), Inches(0.55))
start_circle.fill.solid(); start_circle.fill.fore_color.rgb = C_DARK; start_circle.line.fill.background()
add_text(slide10, "START", 5.9, 1.22, 1.0, 0.45, font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def add_flow_arrow(slide, x, y, h=0.3):
    arr = slide.shapes.add_shape(1, Inches(x+0.35), Inches(y), Inches(0.08), Inches(h))
    arr.fill.solid(); arr.fill.fore_color.rgb = C_GRAY; arr.line.fill.background()
    tip = slide.shapes.add_shape(1, Inches(x+0.23), Inches(y+h-0.01), Inches(0.32), Inches(0.18))
    tip.fill.solid(); tip.fill.fore_color.rgb = C_GRAY; tip.line.fill.background()

def add_process_box(slide, x, y, w, h, text, color, text_color=None, font_size=11):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.color.rgb = RGBColor(0x88,0xAA,0x99); box.line.width = Pt(1)
    tc = text_color if text_color else C_WHITE
    add_text(slide, text, x+0.05, y+(h/2)-0.15, w-0.1, 0.35, font_size=font_size, bold=True, color=tc, align=PP_ALIGN.CENTER)
    return box

# Main flow column x center = 6.39 (box x = 5.9, w=0.9 means center at 6.35... let's use w=1.9, x=5.4)
bx = 5.2; bw = 2.4

# Arrow from START
add_flow_arrow(slide10, bx+0.8, 1.75, 0.3)

# Box 1: Collect Sensor Data
add_process_box(slide10, bx, 2.05, bw, 0.52, "📥  Collect Sensor Data", C_BLUE, font_size=11)
add_flow_arrow(slide10, bx+0.8, 2.57, 0.28)

# Box 2: Preprocess Data
add_process_box(slide10, bx, 2.85, bw, 0.52, "⚙️  Preprocess Data", RGBColor(0x8E,0x44,0xAD), font_size=11)
add_flow_arrow(slide10, bx+0.8, 3.37, 0.28)

# Diamond: Data Valid?
diamond_pts = "0,0.3 0.5,0 1.0,0.3 0.5,0.6"
diam = slide10.shapes.add_shape(4, Inches(bx+0.45), Inches(3.65), Inches(1.5), Inches(0.6))
diam.fill.solid(); diam.fill.fore_color.rgb = RGBColor(0xE6,0x7E,0x22)
diam.line.color.rgb = RGBColor(0xAA,0x55,0x00); diam.line.width = Pt(1.5)
add_text(slide10, "Data\nValid?", bx+0.45, 3.66, 1.5, 0.56, font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# No path (loop back) - left arrow
no_line = slide10.shapes.add_shape(1, Inches(bx-1.2), Inches(3.75), Inches(1.2), Inches(0.08))
no_line.fill.solid(); no_line.fill.fore_color.rgb = C_RED; no_line.line.fill.background()
no_line2 = slide10.shapes.add_shape(1, Inches(bx-1.2), Inches(2.1), Inches(0.08), Inches(1.73))
no_line2.fill.solid(); no_line2.fill.fore_color.rgb = C_RED; no_line2.line.fill.background()
no_line3 = slide10.shapes.add_shape(1, Inches(bx-1.2), Inches(2.1), Inches(1.2), Inches(0.08))
no_line3.fill.solid(); no_line3.fill.fore_color.rgb = C_RED; no_line3.line.fill.background()
add_text(slide10, "No", bx-1.7, 3.6, 0.5, 0.35, font_size=10, bold=True, color=C_RED)
add_text(slide10, "Re-collect", bx-1.75, 2.85, 0.7, 0.35, font_size=9, color=C_RED)

# Yes arrow
add_text(slide10, "Yes", bx+0.9, 4.28, 0.5, 0.3, font_size=10, bold=True, color=C_GREEN)
add_flow_arrow(slide10, bx+0.8, 4.25, 0.3)

# Box 3: Run ML Model
add_process_box(slide10, bx, 4.55, bw, 0.52, "🤖  Run ML Model", C_GREEN, font_size=11)
add_flow_arrow(slide10, bx+0.8, 5.07, 0.28)

# Diamond 2: Irrigation Needed?
diam2 = slide10.shapes.add_shape(4, Inches(bx+0.45), Inches(5.35), Inches(1.5), Inches(0.6))
diam2.fill.solid(); diam2.fill.fore_color.rgb = RGBColor(0x14,0x5A,0x32)
diam2.line.color.rgb = C_MINT; diam2.line.width = Pt(1.5)
add_text(slide10, "Irrigation\nNeeded?", bx+0.45, 5.36, 1.5, 0.56, font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# YES branch: right → Irrigation ON
yes_r_line = slide10.shapes.add_shape(1, Inches(bx+2.4), Inches(5.58), Inches(1.3), Inches(0.08))
yes_r_line.fill.solid(); yes_r_line.fill.fore_color.rgb = C_MINT; yes_r_line.line.fill.background()
on_box = slide10.shapes.add_shape(1, Inches(bx+3.7), Inches(5.35), Inches(2.3), Inches(0.55))
on_box.fill.solid(); on_box.fill.fore_color.rgb = C_MINT; on_box.line.fill.background()
add_text(slide10, "✅  Irrigation ON", bx+3.75, 5.39, 2.2, 0.43, font_size=12, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
add_text(slide10, "YES →", bx+2.42, 5.42, 0.9, 0.3, font_size=9, bold=True, color=C_MINT)

# NO branch: left → Irrigation OFF
no_l_line = slide10.shapes.add_shape(1, Inches(bx-1.7), Inches(5.58), Inches(1.25), Inches(0.08))
no_l_line.fill.solid(); no_l_line.fill.fore_color.rgb = C_GRAY; no_l_line.line.fill.background()
off_box = slide10.shapes.add_shape(1, Inches(bx-4.0), Inches(5.35), Inches(2.3), Inches(0.55))
off_box.fill.solid(); off_box.fill.fore_color.rgb = C_GRAY; off_box.line.fill.background()
add_text(slide10, "❌  Irrigation OFF", bx-3.95, 5.39, 2.2, 0.43, font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide10, "← NO", bx-1.65, 5.42, 0.7, 0.3, font_size=9, bold=True, color=C_GRAY)

# Both branches → Log & Visualize
# Right ON box down arrow
on_down = slide10.shapes.add_shape(1, Inches(bx+4.85), Inches(5.9), Inches(0.08), Inches(0.55))
on_down.fill.solid(); on_down.fill.fore_color.rgb = C_MINT; on_down.line.fill.background()
# Left OFF box down arrow
off_down = slide10.shapes.add_shape(1, Inches(bx-2.85), Inches(5.9), Inches(0.08), Inches(0.55))
off_down.fill.solid(); off_down.fill.fore_color.rgb = C_GRAY; off_down.line.fill.background()

# Horizontal join line
join = slide10.shapes.add_shape(1, Inches(bx-2.85), Inches(6.45), Inches(7.78), Inches(0.08))
join.fill.solid(); join.fill.fore_color.rgb = C_GRAY; join.line.fill.background()

# Down to log box
log_down = slide10.shapes.add_shape(1, Inches(bx+0.8), Inches(6.45), Inches(0.08), Inches(0.3))
log_down.fill.solid(); log_down.fill.fore_color.rgb = C_GRAY; log_down.line.fill.background()

add_process_box(slide10, bx, 6.75, bw, 0.45, "📊  Log Result & Update Dashboard", RGBColor(0x1A,0xBC,0x9C), font_size=10)

# Input labels on left side
input_labels = [
    (0.35, 2.22, "Soil Moisture"),
    (0.35, 2.7, "Temperature"),
    (0.35, 3.18, "Humidity"),
    (0.35, 3.66, "Rainfall"),
]
inp_bg = slide10.shapes.add_shape(1, Inches(0.25), Inches(1.9), Inches(2.2), Inches(2.6))
inp_bg.fill.solid(); inp_bg.fill.fore_color.rgb = RGBColor(0xE8,0xF5,0xFE)
inp_bg.line.color.rgb = C_BLUE; inp_bg.line.width = Pt(1)
add_text(slide10, "INPUT PARAMETERS", 0.3, 1.93, 2.1, 0.35, font_size=9, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
for lx, ly, lt in input_labels:
    dot = slide10.shapes.add_shape(9, Inches(lx+0.1), Inches(ly), Inches(0.15), Inches(0.15))
    dot.fill.solid(); dot.fill.fore_color.rgb = C_BLUE; dot.line.fill.background()
    add_text(slide10, lt, lx+0.35, ly-0.05, 1.7, 0.35, font_size=10, color=C_DARK_TEXT)

# Arrow from input to data collection box
inp_arr = slide10.shapes.add_shape(1, Inches(2.45), Inches(2.28), Inches(2.75), Inches(0.08))
inp_arr.fill.solid(); inp_arr.fill.fore_color.rgb = C_BLUE; inp_arr.line.fill.background()
add_text(slide10, "→", 4.85, 2.12, 0.4, 0.38, font_size=16, bold=True, color=C_BLUE)

# ─────────────────────────────────────────────
# SLIDE 11: Conclusion
# ─────────────────────────────────────────────
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide11, C_DARK)

left_band = slide11.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
left_band.fill.solid(); left_band.fill.fore_color.rgb = C_GREEN; left_band.line.fill.background()

bot = slide11.shapes.add_shape(1, Inches(0), Inches(6.8), Inches(13.33), Inches(0.7))
bot.fill.solid(); bot.fill.fore_color.rgb = C_GREEN; bot.line.fill.background()

add_text(slide11, "Conclusion & Expected Impact", 0.8, 0.7, 11, 0.75, font_size=30, bold=True, color=C_WHITE)
add_text(slide11, "Intelligent Irrigation System Using Machine Learning", 0.8, 1.45, 11, 0.5, font_size=14, color=C_MINT, italic=True)

outcomes = [
    ("💧", "30–50%", "Reduction in water consumption through precision ML-based irrigation"),
    ("⚡", "25%+", "Reduction in electricity costs via optimized pump scheduling"),
    ("📈", "95%+", "Prediction accuracy for irrigation ON/OFF classification"),
    ("🌱", "Better", "Crop yields through optimal soil moisture maintenance"),
]

for i, (icon, val, desc) in enumerate(outcomes):
    cx = 0.8 + i * 3.1
    card = slide11.shapes.add_shape(1, Inches(cx), Inches(2.2), Inches(2.8), Inches(2.0))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x1E,0x3D,0x2D); card.line.fill.background()
    add_text(slide11, icon, cx+0.1, 2.25, 2.6, 0.55, font_size=26, align=PP_ALIGN.CENTER)
    add_text(slide11, val, cx+0.1, 2.78, 2.6, 0.6, font_size=22, bold=True, color=C_MINT, align=PP_ALIGN.CENTER)
    add_text(slide11, desc, cx+0.1, 3.38, 2.6, 0.7, font_size=9.5, color=RGBColor(0xCC,0xEE,0xDD), align=PP_ALIGN.CENTER)

add_text(slide11, "Key Takeaways:", 0.8, 4.5, 11, 0.45, font_size=14, bold=True, color=C_MINT)
takeaways = "This project demonstrates how Machine Learning can transform traditional agriculture by enabling data-driven, real-time irrigation decisions. The system reduces resource waste, lowers operational costs, and improves crop productivity — making smart irrigation accessible to modern farms."
add_text(slide11, takeaways, 0.8, 5.0, 11.5, 1.3, font_size=12, color=RGBColor(0xCC,0xEE,0xDD))

# Save
out_path = "Intelligent_Irrigation_System_ML.pptx"
prs.save(out_path)
print(f"Presentation saved successfully to: {out_path}")
print(f"Saved to {out_path}")